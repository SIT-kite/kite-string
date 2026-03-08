from __future__ import annotations

import argparse
import os
import shutil
import subprocess
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, List, Tuple

from dotenv import load_dotenv


load_dotenv(Path(__file__).resolve().parents[1] / '.env')

OFFICE_DOCUMENT_EXTENSIONS = {
    'doc',
    'docx',
    'ppt',
    'pptx',
    'xls',
    'xlsx',
}

DIRECT_PARSE_EXTENSIONS = {
    'docx',
    'pptx',
    'xls',
    'xlsx',
}

CONTENT_TYPE_MAPPING = {
    'pdf': 'application/pdf',
    'doc': 'application/msword',
    'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    'ppt': 'application/vnd.ms-powerpoint',
    'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    'xls': 'application/vnd.ms-excel',
    'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
}


class UnsupportedAttachmentError(ValueError):
    pass


class MissingParserDependencyError(RuntimeError):
    pass


@dataclass(frozen=True)
class ExtractedDocument:
    content: str
    content_type: str
    parser: str


def split_file_name(file_name: str) -> tuple[str, str]:
    position = file_name.rfind('.')
    if position == -1:
        return file_name, ''
    return file_name[:position], file_name[position + 1:]


def normalize_extension(extension: str | None) -> str:
    if not extension:
        return ''
    return extension.lower().lstrip('.')


def command_exists(command: str) -> bool:
    return shutil.which(command) is not None


def join_text_blocks(blocks: Iterable[str]) -> str:
    return '\n\n'.join(block.strip() for block in blocks if block and block.strip()).strip()


class DocExtractor:
    def __init__(self, path: str = ''):
        self._path = path


class PdfExtractor(DocExtractor):
    def page2png(self, output_directory: str | None = None) -> List[str]:
        if not output_directory:
            name, _ = split_file_name(self._path)
            output_directory = name
            os.makedirs(output_directory, exist_ok=True)

        from wand.image import Image

        pdf2img_obj = Image(filename=self._path, resolution=300)
        png_objs = pdf2img_obj.convert('png')

        output_files = []
        for index, img in enumerate(png_objs.sequence):
            current_page = Image(image=img)
            file_name = f'{output_directory}/{index}.png'

            with open(file_name, 'wb') as out_image:
                out_image.write(current_page.make_blob('png'))
            output_files.append(file_name)

        return output_files

    def text(self) -> List[Tuple[int, str]]:
        from pdfminer.converter import PDFPageAggregator
        from pdfminer.layout import LAParams, LTTextBoxHorizontal
        from pdfminer.pdfinterp import PDFPageInterpreter, PDFResourceManager
        from pdfminer.pdfparser import PDFDocument, PDFParser

        with open(self._path, 'rb') as pdf_file:
            doc = PDFDocument()
            parser = PDFParser(pdf_file)
            parser.set_document(doc)
            doc.set_parser(parser)
            doc.initialize()
            if not doc.is_extractable:
                raise RuntimeError(f'pdf text extraction is not allowed for {self._path}')

            rsrcmgr = PDFResourceManager()
            laparams = LAParams()
            device = PDFPageAggregator(rsrcmgr, laparams=laparams)
            interpreter = PDFPageInterpreter(rsrcmgr, device)

            result = []
            for index, page in enumerate(doc.get_pages()):
                interpreter.process_page(page)
                layout = device.get_result()

                current_page_text = ''.join(
                    element.get_text().strip()
                    for element in layout
                    if isinstance(element, LTTextBoxHorizontal)
                )
                result.append((index, current_page_text))

        return result


class DocxExtractor(DocExtractor):
    def text(self) -> List[Tuple[int, str]]:
        from docx import Document

        document = Document(self._path)
        blocks: list[str] = []
        blocks.extend(paragraph.text for paragraph in document.paragraphs)
        for table in document.tables:
            for row in table.rows:
                blocks.append('\t'.join(cell.text for cell in row.cells))
        return [(0, join_text_blocks(blocks))]


class XlsxExtractor(DocExtractor):
    def text(self) -> List[Tuple[int, str]]:
        from openpyxl import load_workbook

        workbook = load_workbook(self._path, read_only=True, data_only=True)
        pages: list[Tuple[int, str]] = []
        for index, sheet in enumerate(workbook.worksheets):
            rows = []
            for row in sheet.iter_rows(values_only=True):
                values = [str(value).strip() for value in row if value is not None and str(value).strip()]
                if values:
                    rows.append('\t'.join(values))
            pages.append((index, join_text_blocks([sheet.title, *rows])))
        return pages


class XlsExtractor(DocExtractor):
    def text(self) -> List[Tuple[int, str]]:
        import xlrd

        workbook = xlrd.open_workbook(self._path)
        pages: list[Tuple[int, str]] = []
        for index in range(workbook.nsheets):
            sheet = workbook.sheet_by_index(index)
            rows = []
            for row_index in range(sheet.nrows):
                values = []
                for cell in sheet.row_values(row_index):
                    text = str(cell).strip()
                    if text:
                        values.append(text)
                if values:
                    rows.append('\t'.join(values))
            pages.append((index, join_text_blocks([sheet.name, *rows])))
        return pages


class PptxExtractor(DocExtractor):
    def text(self) -> List[Tuple[int, str]]:
        from pptx import Presentation

        presentation = Presentation(self._path)
        pages: list[Tuple[int, str]] = []
        for index, slide in enumerate(presentation.slides):
            blocks = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    blocks.append(shape.text)
            pages.append((index, join_text_blocks(blocks)))
        return pages


class DocBinaryExtractor(DocExtractor):
    def text(self) -> List[Tuple[int, str]]:
        if command_exists('antiword'):
            try:
                result = subprocess.run(
                    ['antiword', self._path],
                    check=True,
                    capture_output=True,
                    text=True,
                    timeout=120,
                )
                return [(0, result.stdout.strip())]
            except subprocess.CalledProcessError:
                pass

        if command_exists('catdoc'):
            result = subprocess.run(
                ['catdoc', self._path],
                check=True,
                capture_output=True,
                text=True,
                timeout=120,
            )
            return [(0, result.stdout.strip())]

        raise MissingParserDependencyError('antiword, catdoc or libreoffice is required to parse this file type')


class OfficeExtractor(DocExtractor):
    @staticmethod
    def convert_file_to_pdf(document_file: str, output_directory: str | None = None) -> str:
        document_path = os.path.abspath(document_file)
        if output_directory is None:
            output_directory = os.path.dirname(document_path)
        output_directory = os.path.abspath(output_directory)
        os.makedirs(output_directory, exist_ok=True)

        office_bin = os.getenv('LIBREOFFICE_BIN', 'libreoffice')
        if not command_exists(office_bin):
            raise MissingParserDependencyError(
                f'{office_bin} is required to parse this file type'
            )

        command = [
            office_bin,
            '--headless',
            '--convert-to',
            'pdf',
            document_path,
            '--outdir',
            output_directory,
        ]
        subprocess.run(command, check=True, timeout=120)

        file_title, _ = split_file_name(os.path.basename(document_path))
        return os.path.join(output_directory, file_title + '.pdf')

    def convert_to_pdf(self, output_directory: str | None = None) -> str:
        return self.convert_file_to_pdf(self._path, output_directory)

    def convert2png(self, output_directory: str | None = None):
        pdf_path = self.convert_to_pdf()
        PdfExtractor(pdf_path).page2png(output_directory)

    def text(self) -> List[Tuple[int, str]]:
        with tempfile.TemporaryDirectory(prefix='kite-attachment-') as output_directory:
            pdf_path = self.convert_to_pdf(output_directory)
            return PdfExtractor(pdf_path).text()


class WordExtractor(OfficeExtractor):
    pass


DIRECT_EXTRACTOR_MAPPING = {
    'doc': DocBinaryExtractor,
    'docx': DocxExtractor,
    'pptx': PptxExtractor,
    'xls': XlsExtractor,
    'xlsx': XlsxExtractor,
}


def extract_text(path: str, extension: str | None = None) -> ExtractedDocument:
    normalized_extension = normalize_extension(extension or Path(path).suffix)

    if normalized_extension == 'pdf':
        pages = PdfExtractor(path).text()
        parser = 'pdfminer3k'
    elif normalized_extension in DIRECT_EXTRACTOR_MAPPING:
        pages = DIRECT_EXTRACTOR_MAPPING[normalized_extension](path).text()
        parser = DIRECT_EXTRACTOR_MAPPING[normalized_extension].__name__.lower()
    elif normalized_extension in OFFICE_DOCUMENT_EXTENSIONS:
        pages = OfficeExtractor(path).text()
        parser = 'libreoffice+pdfminer3k'
    else:
        raise UnsupportedAttachmentError(f'unsupported extension: {normalized_extension or "<empty>"}')

    content = '\n\n'.join(text.strip() for _, text in pages if text and text.strip()).strip()
    content_type = CONTENT_TYPE_MAPPING.get(normalized_extension, 'application/octet-stream')
    return ExtractedDocument(content=content, content_type=content_type, parser=parser)


def main():
    parser = argparse.ArgumentParser(description='Extract text from a document.')
    parser.add_argument('path', help='Document path on disk')
    parser.add_argument('--ext', default='', help='Override file extension')
    arguments = parser.parse_args()

    extracted = extract_text(arguments.path, arguments.ext)
    print(extracted.content)


if __name__ == '__main__':
    main()
